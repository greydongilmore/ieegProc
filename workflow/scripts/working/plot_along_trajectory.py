#!/usr/bin/env python3
# -*- coding: utf-8 -*-


import numpy as np
import nibabel as nib
from nilearn import image
import matplotlib.pyplot as plt
import glob
import pandas as pd
import os
from matplotlib.patches import Rectangle
from PIL import Image
import SimpleITK as sitk

os.chdir(r'/home/greydon/Documents/GitHub/seeg2bids-pipeline/workflow/scripts/working')
from helpers import electrodeModels,determineFCSVCoordSystem,determine_groups,normDir,norm_vec

class IndexTracker:
	def __init__(self, ax, img_data,points=None, rotate_img=False, rotate_points=False, title=None):
		self.ax = ax
		self.scatter=None
		self.points=None
		self.title=title
		no_index= True
		if rotate_img:
			self.img_data = np.fliplr(np.rot90(img_data.copy(),3))
			#if self.points is not None:
			#	self.points = np.rot90(points, k=1)
		else:
			self.img_data = img_data.copy()
		
		rows, cols, self.slices = self.img_data.shape
		self.ind = self.slices//2
		if points is not None:
			self.points=points.copy()
			if rotate_points:
				self.points[:,[0,1]]=self.points[:,[1,0]]
			while no_index==True:
				if any(self.points[:,2]==self.ind):
					no_index=False
					point_plot=np.vstack([np.mean(self.points[(self.points[:,2]==self.ind)*(self.points[:,3]==x),:2],0) for x in np.unique(self.points[self.points[:,2]==self.ind,3])])
					self.scatter,=ax.plot(point_plot[:,1],point_plot[:,0], marker="o", markersize=12, c = 'yellow', fillstyle='none', markeredgewidth=1, linestyle = 'None')
				else:
					self.ind+=1
				
		self.im = ax.imshow(self.img_data[:, :, self.ind], origin='lower')
		self.update()
	
	def on_scroll(self, event):
		print("%s %s" % (event.button, event.step))
		if event.button == 'up':
			self.ind = (self.ind + 1) % self.slices
		else:
			self.ind = (self.ind - 1) % self.slices
		self.update()
	
	def update(self):
		
		if self.points is not None:
			if any(self.points[:,2]==self.ind):
				point_plot=np.vstack([np.mean(self.points[(self.points[:,2]==self.ind)*(self.points[:,3]==x),:2],0) for x in np.unique(self.points[self.points[:,2]==self.ind,3])])
				self.scatter.set_xdata(point_plot[:,1])
				self.scatter.set_ydata(point_plot[:,0])
		self.im.set_data(self.img_data[:, :, self.ind])
		
		if self.title is not None:
			plot_title=self.title
		else:
			plot_title='slice %s' % self.ind
		
		self.ax.set_title(plot_title, fontdict={'fontsize': 18, 'fontweight': 'bold'})
		self.ax.tick_params(axis='x', labelsize=14)
		self.ax.tick_params(axis='y', labelsize=14)
		self.im.axes.figure.canvas.draw()
		

def determineImageThreshold( img_data):
	hist_y, hist_x = np.histogram(img_data.flatten(), bins=256)
	hist_x = hist_x[0:-1]
	
	cumHist_y = np.cumsum(hist_y.astype(float))/np.prod(np.array(img_data.shape))
	# The background should contain half of the voxels
	minThreshold_byCount = hist_x[np.where(cumHist_y > 0.90)[0][0]]
	
	hist_diff = np.diff(hist_y)
	hist_diff_zc = np.where(np.diff(np.sign(hist_diff)) == 2)[0].flatten()
	minThreshold = hist_x[hist_diff_zc[hist_x[hist_diff_zc] > (minThreshold_byCount)][0]]
	print(f"first maxima after soft tissue found: {minThreshold}")
	
	
	return minThreshold

def bounding_box(seg):
	x = np.any(np.any(seg, axis=0), axis=1)
	y = np.any(np.any(seg, axis=1), axis=1)
	z = np.any(np.any(seg, axis=1), axis=0)
	ymin, ymax = np.where(y)[0][[0, -1]]
	xmin, xmax = np.where(x)[0][[0, -1]]
	zmin, zmax = np.where(z)[0][[0, -1]]
	bbox = np.array([ymin,ymax,xmin,xmax,zmin,zmax])
	return bbox

def bbox2(img):
	rows = np.any(img, axis=(1, 2))
	cols = np.any(img, axis=(0, 2))
	z = np.any(img, axis=(0, 1))
	
	ymin, ymax = np.where(rows)[0][[0, -1]]
	xmin, xmax = np.where(cols)[0][[0, -1]]
	zmin, zmax = np.where(z)[0][[0, -1]]
	return img[ymin:ymax+1, xmin:xmax+1, zmin:zmax+1]

def nib_to_sitk(data, affine):
	FLIP_XY = np.diag((-1, -1, 1))
	origin = np.dot(FLIP_XY, affine[:3, 3]).astype(np.float64)
	RZS = affine[:3, :3]
	spacing = np.sqrt(np.sum(RZS * RZS, axis=0))
	R = RZS / spacing
	direction = np.dot(FLIP_XY, R).flatten()
	image = sitk.GetImageFromArray(data.transpose())
	image.SetOrigin(origin)
	image.SetSpacing(spacing)
	image.SetDirection(direction)
	return image


#%%


debug = False
if debug:
	class dotdict(dict):
		"""dot.notation access to dictionary attributes"""
		__getattr__ = dict.get
		__setattr__ = dict.__setitem__
		__delattr__ = dict.__delitem__

	class Namespace:
		def __init__(self, **kwargs):
			self.__dict__.update(kwargs)
	
	isub = 'sub-P108'
	data_dir = r'/home/greydon/Documents/data/SEEG/derivatives'
	
	input = dotdict({
			'actual_fcsv': f'{data_dir}/seega_scenes/{isub}/*actual.fcsv',
			'planned_fcsv': f'{data_dir}/seega_scenes/{isub}/*planned.fcsv',
			'shopping_list': f'{data_dir}/seega_scenes/{isub}/*shopping_list.xlsx',
			't1w_fname': f'{data_dir}/seega_scenes/{isub}/{isub}*_acq-contrast_T1w.nii.gz',
			'ct_fname': f'{data_dir}/seega_scenes/{isub}/{isub}*_ct.nii.gz',
	 })
	
	output = dotdict({
			'output_dir': f'{data_dir}/seega_scenes/{isub}/traj_snapshots',
	 })
	
	snakemake = Namespace(input=input,output=output)


if glob.glob(snakemake.input.actual_fcsv):
	subj=os.path.basename(glob.glob(snakemake.input.t1w_fname)[0]).split('_')[0]
	
	if not os.path.exists(snakemake.output.output_dir):
		os.makedirs(snakemake.output.output_dir)
	
	df_elec_raw = pd.read_excel(glob.glob(snakemake.input.shopping_list)[0],header=None)
	df_elec=df_elec_raw.iloc[4:,:].reset_index(drop=True)
	df_elec.columns=df_elec_raw.iloc[3]
	df_elec=df_elec.iloc[0:df_elec.iloc[:,1].isnull().idxmax()]
	
	actual_version_num, actual_head_info=determineFCSVCoordSystem(glob.glob(snakemake.input.actual_fcsv)[0])
	actual_data = pd.read_csv(glob.glob(snakemake.input.actual_fcsv)[0], skiprows=3, header=None)
	actual_data = actual_data.rename(columns=actual_head_info).reset_index(drop=True)
	groupsActual, actual_all = determine_groups(actual_data['label'].values)
	
	planned_version_num, planned_head_info=determineFCSVCoordSystem(glob.glob(snakemake.input.planned_fcsv)[0])
	planned_data = pd.read_csv(glob.glob(snakemake.input.planned_fcsv)[0], skiprows=3, header=None)
	planned_data = planned_data.rename(columns=planned_head_info).reset_index(drop=True)
	groupsPlanned, planned_all = determine_groups(planned_data['label'].values)
	
	for igroup in set(groupsActual).intersection(groupsPlanned):
		e_info=None
		if any(igroup in x for x in df_elec['Target'].values):
			group_idx=[i for i,x in enumerate(df_elec['Target'].values.tolist()) if igroup in x][0]
			elec_str=df_elec.loc[group_idx,'Electrode']
			if isinstance(elec_str,int):
				if 'mm' in df_elec.loc[group_idx,'Suggested']:
					elec_str=str(elec_str)+' mm'
				elif 'contact' in df_elec.loc[group_idx,'Suggested']:
					elec_str=str(elec_str)+' contact'
			
			elecIdx=[i for i,x in enumerate(list(electrodeModels)) if elec_str == x]
			if elecIdx:
				e_info=electrodeModels[list(electrodeModels)[elecIdx[0]]]
			
			actual_idx=[i for i,x in enumerate(actual_data['label'].values) if x.startswith(igroup)]
			actual_target=actual_data.loc[actual_idx,['x','y','z']].values.astype(float)[0]
			actual_entry=actual_data.loc[actual_idx,['x','y','z']].values.astype(float)[1]
			
			planned_idx=[i for i,x in enumerate(planned_data['label'].values) if x.startswith(igroup)]
			planned_target=planned_data.loc[planned_idx,['x','y','z']].values.astype(float)[0]
			planned_entry=planned_data.loc[planned_idx,['x','y','z']].values.astype(float)[1]
			
			direction=actual_target-actual_entry
			xAxis = normDir(np.cross([1,0,0],direction))
			yAxis = normDir(np.cross(direction,xAxis))
			
			for iplot in ('coronal','sagittal'):
				trajMatrix = np.identity(3)
				if iplot=='axial':
					trajMatrix[0:3,0]=normDir(np.cross(direction,yAxis))
					trajMatrix[0:3,1]=normDir(yAxis)
					trajMatrix[0:3,2]=normDir(direction)
				elif iplot =='sagittal':
					trajMatrix[0:3,0]=normDir(np.cross(direction,yAxis))
					trajMatrix[0:3,1]=normDir(direction)
					trajMatrix[0:3,2]=normDir(yAxis)
				elif iplot =='coronal':
					trajMatrix[0:3,0]=normDir(-yAxis)
					trajMatrix[0:3,1]=normDir(direction)
					trajMatrix[0:3,2]=normDir(np.cross(direction,-yAxis))
				
				t1w_obj=nib.load(glob.glob(snakemake.input.t1w_fname)[0])
				ct_obj=nib.load(glob.glob(snakemake.input.ct_fname)[0])
				
				t1w_reslice = image.resample_img(t1w_obj, target_affine=trajMatrix, interpolation='linear')
				ct_reslice = image.resample_to_img(ct_obj, t1w_reslice, interpolation='linear')
				#ct_thres=ct_reslice.get_fdata()*(ct_reslice.get_fdata()>1200)
				
				t1w_reslice_sitk=nib_to_sitk(t1w_reslice.get_fdata(),t1w_reslice.affine)
				voxsize = t1w_reslice_sitk.GetSpacing()
				
				minThreshold=determineImageThreshold(t1w_reslice.get_fdata())
				thresh_img = t1w_reslice_sitk > minThreshold
				
				box = bounding_box(t1w_reslice.get_fdata())
				t1w_reslice_cut = t1w_reslice.get_fdata()[box[0]:box[1]+1,box[2]:box[3]+1,box[4]:box[5]+1]
				ct_reslice_cut = ct_reslice.get_fdata()[box[0]:box[1]+1,box[2]:box[3]+1,box[4]:box[5]+1]
				
				
				plt_img=sitk.GetArrayFromImage(thresh_img)
				
				fig, ax = plt.subplots(1, 1)
				tracker = IndexTracker(ax, ct_reslice.get_fdata(), points=None,rotate_img=True, rotate_points=True)
				fig.canvas.mpl_connect('scroll_event', tracker.on_scroll)
				plt.show()
				
				planned_target_trans=(np.linalg.inv(t1w_reslice.affine)@(np.append(planned_target,1))).astype(int)
				planned_entry_trans=(np.linalg.inv(t1w_reslice.affine)@(np.append(planned_entry,1))).astype(int)
				actual_target_trans=(np.linalg.inv(t1w_reslice.affine)@(np.append(actual_target,1))).astype(int)
				actual_entry_trans=(np.linalg.inv(t1w_reslice.affine)@(np.append(actual_entry,1))).astype(int)
				
				center_dcm=((actual_entry_trans+actual_target_trans)/2).astype(int)
				center_planned_dcm=((planned_entry_trans+planned_target_trans)/2).astype(int)
				
				yshift=3
				planned_coords=[(planned_target_trans[0],planned_entry_trans[0]), (planned_target_trans[1]+yshift,planned_entry_trans[1]+yshift)]
				
				offset=0
				x_shift=3
				if iplot=='axial':
					rec_coords=[planned_entry_trans[2]-20,planned_entry_trans[1],(planned_target_trans[2]-planned_entry_trans[2])+20,2]
				elif iplot=='sagittal' or iplot=='coronal':
					rec_coords=[planned_target_trans[2]-x_shift,planned_target_trans[1],x_shift,(planned_target_trans[1]-planned_entry_trans[1])+20]
				
				#plt.ioff()
				fig = plt.figure(figsize=(12,12))
				ax = fig.add_subplot(111)
				ax.imshow(ct_reslice.get_fdata()[:,:,actual_target_trans[2]], cmap='gray',alpha=1,origin="lower",vmax=2500)
				plt.plot(planned_coords[1], planned_coords[0], color="blue", linewidth=1)
				
				ax.imshow(np.rot90(np.fliplr(t1w_reslice_cut[actual_target_trans[0]-box[0],:,:]),2), cmap='gray',alpha=1,origin="lower",vmax=2500)
				ax.imshow(np.rot90(np.fliplr(ct_reslice_cut[actual_target_trans[0]-box[0],:,:]),2),cmap='gray',alpha=.7,origin="lower",vmin=1200)
				plt.gca().add_patch(Rectangle((rec_coords[0],rec_coords[1]),rec_coords[2],rec_coords[3],edgecolor='deepskyblue',facecolor='none',lw=1.5,alpha=1))
				
				if e_info is not None:
					y_index=rec_coords[1]+e_info['encapsultation']+offset
					x_index=rec_coords[0]
					for icontact in range(e_info['num_contacts']):
						ax.plot((x_index,x_index),(y_index,y_index+e_info['contact_size']),color='fuchsia',linewidth=1.5)
						ax.plot((x_index+x_shift,x_index+x_shift),(y_index,y_index+e_info['contact_size']),color='fuchsia',linewidth=1.5)
						if icontact < (e_info['num_contacts']-1):
							y_index+=(e_info['contact_spacing'][icontact]+e_info['contact_size'])
					
				ax.axis('off')
				plt.tight_layout()
				plt.savefig(os.path.join(snakemake.output.output_dir, f"{subj}_desc-{iplot}_traj-{igroup}_qc.svg"),transparent=True,dpi=350)
				plt.ion()

#%%
import svgutils.transform as sg
import cairosvg
from pptx import Presentation
import io
from pptx.util import Inches,Pt
import vtk


prs = Presentation(r'/home/greydon/Documents/data/SEEG/derivatives/seega_scenes/sub-P076/Corkery_Danyelle_2021-11-25_maps.pptx')
for iphoto in os.listdir(snakemake.output.output_dir):
	label=iphoto.split('_')[2].split('-')[-1]
	slide_idx=[i for i,x in enumerate(prs.slides) if label in x.name]
	if slide_idx:
		slide = prs.slides[slide_idx[0]]
		
		image_fname=os.path.join(snakemake.output.output_dir,iphoto)
		png_data = cairosvg.svg2png(url=image_fname, write_to=os.path.splitext(image_fname)[0]+'.png', dpi=350)
		
		top=Inches(prs.slide_height.inches-Pt(400).inches)/2
		pic = slide.shapes.add_picture(os.path.splitext(image_fname)[0]+'.png',0,top,height=Pt(440))
		prs.save(r'/home/greydon/Documents/data/SEEG/derivatives/seega_scenes/sub-P076/Corkery_Danyelle_2021-11-25_maps.pptx')



entryTargetDirection = actual_entry - actual_target
vtk.vtkMath().Normalize(entryTargetDirection)
superiorInferiorDirection = np.array([0,0,1])

ang_rad = np.arccos(vtk.vtkMath().Dot(entryTargetDirection, superiorInferiorDirection))
ang_deg = np.rad2deg(ang_rad)

cross = np.zeros(3)
vtk.vtkMath().Cross(entryTargetDirection, superiorInferiorDirection, cross)

if vtk.vtkMath().Dot(cross,superiorInferiorDirection) >= 0:
	ang_deg = -1 * ang_deg

vtkTransform = vtk.vtkTransform()
vtkTransform.Translate(actual_target)
vtkTransform.RotateWXYZ(rollAngle, entryTargetDirection[0], entryTargetDirection[1], entryTargetDirection[2])
vtkTransform.RotateWXYZ(ang_deg, cross[0], cross[1], cross[2])

outputTransform.SetAndObserveTransformToParent(vtkTransform)
